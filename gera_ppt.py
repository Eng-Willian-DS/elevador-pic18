"""
gera_ppt.py — Gerador de apresentação: Central para Elevador
ELE1012 — Microcontroladores — UTFPR Medianeira
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Cm
import copy

# ─── PALETA DE CORES ─────────────────────────────────────────────
BG_DARK    = RGBColor(0x0D, 0x11, 0x17)   # fundo principal
BG_CARD    = RGBColor(0x16, 0x1B, 0x22)   # fundo de cards/código
ACCENT     = RGBColor(0xF9, 0xA8, 0x25)   # amarelo UTFPR
ACCENT2    = RGBColor(0x58, 0xA6, 0xFF)   # azul destaque
GREEN      = RGBColor(0x3F, 0xB9, 0x50)   # verde (código)
WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
GRAY       = RGBColor(0x8B, 0x94, 0x9E)
RED        = RGBColor(0xFF, 0x5C, 0x5C)

# Cor por integrante
COR_WILLIAN  = RGBColor(0x58, 0xA6, 0xFF)   # azul
COR_KAROLINE = RGBColor(0xBB, 0x87, 0xFF)   # roxo
COR_ANA      = RGBColor(0x3F, 0xB9, 0x50)   # verde
COR_ERICK    = RGBColor(0xFF, 0x7B, 0x72)   # vermelho/laranja

W  = Inches(13.33)   # largura slide widescreen
H  = Inches(7.5)     # altura slide widescreen

prs = Presentation()
prs.slide_width  = W
prs.slide_height = H

BLANK = prs.slide_layouts[6]  # layout em branco

# ─── HELPERS ─────────────────────────────────────────────────────

def add_slide():
    s = prs.slides.add_slide(BLANK)
    # fundo preto
    bg = s.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = BG_DARK
    return s

def box(slide, x, y, w, h, color=None, alpha=None):
    shape = slide.shapes.add_shape(1, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.line.fill.background()
    if color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = color
    else:
        shape.fill.background()
    return shape

def txt(slide, text, x, y, w, h,
        size=18, bold=False, color=WHITE, align=PP_ALIGN.LEFT,
        wrap=True, italic=False):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tb.word_wrap = wrap
    tf = tb.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    run.font.name = "Segoe UI"
    return tb

def txt_code(slide, lines, x, y, w, h, size=11):
    """Bloco de código com fundo escuro."""
    card = box(slide, x, y, w, h, color=BG_CARD)
    tb = slide.shapes.add_textbox(
        Inches(x + 0.15), Inches(y + 0.1),
        Inches(w - 0.3), Inches(h - 0.2))
    tb.word_wrap = False
    tf = tb.text_frame
    tf.word_wrap = False
    for i, line in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        run = p.add_run()
        run.text = line
        run.font.size = Pt(size)
        run.font.name = "Consolas"
        run.font.color.rgb = GREEN
    return card

def header_bar(slide, label, color, parte_num=None):
    """Barra superior colorida com label do integrante."""
    bar = box(slide, 0, 0, 13.33, 0.55, color=color)
    text = f"  {label}"
    if parte_num:
        text = f"  PARTE {parte_num} — {label}"
    txt(slide, text, 0.1, 0.02, 13, 0.5,
        size=16, bold=True, color=BG_DARK, align=PP_ALIGN.LEFT)

def slide_title(slide, title, subtitle=None, accent=ACCENT):
    """Título grande centralizado."""
    txt(slide, title, 0.5, 1.5, 12.33, 1.5,
        size=40, bold=True, color=accent, align=PP_ALIGN.CENTER)
    if subtitle:
        txt(slide, subtitle, 0.5, 3.1, 12.33, 1.0,
            size=22, color=GRAY, align=PP_ALIGN.CENTER)

def divider(slide, y, color=ACCENT, width=12.33, x=0.5):
    bar = box(slide, x, y, width, 0.04, color=color)
    return bar

def bullet_block(slide, items, x, y, w, h, size=16, color=WHITE, icon="▸"):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tb.word_wrap = True
    tf = tb.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.space_before = Pt(4)
        run = p.add_run()
        run.text = f"{icon}  {item}"
        run.font.size = Pt(size)
        run.font.name = "Segoe UI"
        run.font.color.rgb = color

def label_box(slide, label, x, y, w=2.5, h=0.4, color=ACCENT):
    b = box(slide, x, y, w, h, color=color)
    txt(slide, label, x+0.05, y+0.03, w-0.1, h-0.06,
        size=13, bold=True, color=BG_DARK, align=PP_ALIGN.CENTER)

# ═══════════════════════════════════════════════════════════════
#  SLIDE 1 — CAPA
# ═══════════════════════════════════════════════════════════════
s = add_slide()
box(s, 0, 3.1, 13.33, 0.08, color=ACCENT)
box(s, 0, 3.2, 13.33, 0.04, color=ACCENT2)

txt(s, "CENTRAL PARA ELEVADOR", 0.5, 0.6, 12.33, 1.5,
    size=44, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)
txt(s, "Controle de 6 andares com PIC18F4550", 0.5, 2.0, 12.33, 0.8,
    size=22, color=GRAY, align=PP_ALIGN.CENTER)

# Integrantes
membros = [
    ("Willian Douglas dos Santos Clementino", "RA 1869027", COR_WILLIAN),
    ("Karoline Vitória de Freitas Yang",      "RA 2578000", COR_KAROLINE),
    ("Ana Caroline Soares Baptista",           "RA 2486024", COR_ANA),
    ("Erick Rodrigues Villalva",               "RA 2301105", COR_ERICK),
]
xs = [0.3, 3.6, 6.9, 10.0]
for i, ((nome, ra, cor), x) in enumerate(zip(membros, xs)):
    box(s, x, 3.5, 3.0, 2.2, color=BG_CARD)
    box(s, x, 3.5, 3.0, 0.07, color=cor)
    txt(s, nome, x+0.12, 3.6, 2.8, 1.2,
        size=13, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    txt(s, ra, x+0.12, 4.7, 2.8, 0.45,
        size=13, color=cor, align=PP_ALIGN.CENTER, bold=True)

txt(s, "ELE1012 — Microcontroladores  |  Prof. Miyadaira  |  UTFPR Medianeira  |  2026",
    0.5, 6.9, 12.33, 0.5, size=13, color=GRAY, align=PP_ALIGN.CENTER)

# ═══════════════════════════════════════════════════════════════
#  SLIDE 2 — VISÃO GERAL DO PROJETO
# ═══════════════════════════════════════════════════════════════
s = add_slide()
box(s, 0, 0, 13.33, 0.55, color=ACCENT)
txt(s, "  VISÃO GERAL DO PROJETO", 0.1, 0.05, 12, 0.45,
    size=18, bold=True, color=BG_DARK)

txt(s, "Requisitos do trabalho", 0.5, 0.75, 12, 0.6,
    size=26, bold=True, color=ACCENT)
divider(s, 1.4)

requisitos = [
    "6 andares — botões de chamada por andar (RB0–RB5)",
    "Display LCD 2×16 mostrando andar atual, temperatura e sentido (↑ ↓)",
    "Temperatura simulada por potenciômetro — elevador para se temp > 31 °C",
    "Após 2,3 s inativo, retorna automaticamente ao 3º andar",
    "Ao chegar no andar desejado, porta abre e permanece aberta por 1 s",
    "Sinal PWM para o motor: aproximadamente 7 kHz com ciclo ativo de 38%",
]
bullet_block(s, requisitos, 0.5, 1.55, 12.3, 4.5, size=18, icon="✔")

txt(s, "Plataforma: PIC18F4550  |  IDE: MPLAB v8.92 + C18 v3.47 LITE  |  Simulação: PICSimLab",
    0.5, 6.85, 12.33, 0.5, size=13, color=GRAY, align=PP_ALIGN.CENTER)

# ═══════════════════════════════════════════════════════════════
#  SLIDE 3 — DIVISÃO DE CONTEÚDO
# ═══════════════════════════════════════════════════════════════
s = add_slide()
box(s, 0, 0, 13.33, 0.55, color=ACCENT2)
txt(s, "  DIVISÃO DO TRABALHO", 0.1, 0.05, 12, 0.45,
    size=18, bold=True, color=BG_DARK)

txt(s, "Quem apresenta o quê", 0.5, 0.7, 12, 0.6,
    size=26, bold=True, color=WHITE)
divider(s, 1.35, color=ACCENT2)

partes = [
    (COR_WILLIAN,  "WILLIAN",    "PARTE 1",
     ["Hardware: PIC18F4550, config bits, pinagem",
      "Cálculo e configuração do PWM 7 kHz / 38%",
      "Funções: init_hardware(), init_pwm()"]),
    (COR_KAROLINE, "KAROLINE",   "PARTE 2",
     ["Timer0: ISR de 1 ms, cálculo de recarga",
      "ADC: leitura do potenciômetro, mapeamento 0–50 °C",
      "Função: init_timer0(), init_adc(), ler_adc()"]),
    (COR_ANA,      "ANA CAROLINE","PARTE 3",
     ["LCD Miyadaira 4-bit, lcd_inicia(), lcd_posicao()",
      "CGRAM: criação de ↑ e ↓ em carregar_chars()",
      "Botões: ler_botoes() com debounce de 20 ms"]),
    (COR_ERICK,    "ERICK",      "PARTE 4",
     ["Máquina de estados (IDLE→MOVENDO→PORTA→SUPER)",
      "Funções motor_liga/desliga(), porta_abre/fecha()",
      "Display: atualizar_lcd() com temperatura e seta"]),
]

for i, (cor, nome, parte, items) in enumerate(partes):
    x = 0.3 + i * 3.2
    box(s, x, 1.55, 3.0, 4.9, color=BG_CARD)
    box(s, x, 1.55, 3.0, 0.55, color=cor)
    txt(s, parte, x+0.1, 1.57, 2.8, 0.28,
        size=11, bold=True, color=BG_DARK, align=PP_ALIGN.LEFT)
    txt(s, nome, x+0.1, 1.85, 2.8, 0.28,
        size=13, bold=True, color=BG_DARK, align=PP_ALIGN.LEFT)
    bullet_block(s, items, x+0.1, 2.2, 2.8, 4.0, size=13, color=WHITE, icon="›")

# ═══════════════════════════════════════════════════════════════
#  SLIDE 4 — WILLIAN: Hardware + Pinagem
# ═══════════════════════════════════════════════════════════════
s = add_slide()
header_bar(s, "WILLIAN DOUGLAS  —  RA 1869027", COR_WILLIAN, "1")

txt(s, "Hardware — PIC18F4550 e Pinagem", 0.5, 0.7, 12, 0.7,
    size=28, bold=True, color=COR_WILLIAN)
divider(s, 1.45, color=COR_WILLIAN)

# tabela de pinos
pinos = [
    ("RB0–RB5", "PORTB", "Botões dos andares 1–6  (pull-up interno, ativo em LOW)"),
    ("RA0/AN0", "PORTA", "Potenciômetro — simula temperatura 0–50 °C  (ADC)"),
    ("RC2/CCP1","PORTC", "Saída PWM do motor  (~7 kHz, 38% duty cycle)"),
    ("RD0",     "PORTD", "Motor IN1 — define direção de giro (via L293D)"),
    ("RD1",     "PORTD", "Motor IN2 — define direção de giro (via L293D)"),
    ("RD2",     "PORTD", "LED porta  (1 = aberta, 0 = fechada)"),
    ("RD4–RD7", "PORTD", "LCD data bus 4 bits  (biblioteca Miyadaira)"),
    ("RE0–RE2", "PORTE", "LCD: RW (RE0) · Enable (RE1) · RS (RE2)"),
]

# cabeçalho tabela
box(s, 0.3, 1.55, 2.0, 0.38, color=ACCENT)
box(s, 2.3, 1.55, 1.8, 0.38, color=ACCENT)
box(s, 4.1, 1.55, 8.8, 0.38, color=ACCENT)
for label, col in [("Pino", 0.3), ("Porta", 2.3), ("Função", 4.1)]:
    txt(s, label, col+0.1, 1.58, 2.5, 0.35,
        size=13, bold=True, color=BG_DARK, align=PP_ALIGN.LEFT)

for i, (pino, porta, funcao) in enumerate(pinos):
    y = 1.93 + i * 0.55
    bg = BG_CARD if i % 2 == 0 else RGBColor(0x1E, 0x25, 0x33)
    box(s, 0.3, y, 2.0, 0.52, color=bg)
    box(s, 2.3, y, 1.8, 0.52, color=bg)
    box(s, 4.1, y, 8.8, 0.52, color=bg)
    txt(s, pino,   0.4,  y+0.1, 2.0, 0.4, size=13, bold=True,  color=GREEN)
    txt(s, porta,  2.4,  y+0.1, 1.6, 0.4, size=13, color=GRAY)
    txt(s, funcao, 4.2,  y+0.1, 8.6, 0.4, size=13, color=WHITE)

# ═══════════════════════════════════════════════════════════════
#  SLIDE 5 — WILLIAN: Config bits + init_hardware()
# ═══════════════════════════════════════════════════════════════
s = add_slide()
header_bar(s, "WILLIAN DOUGLAS  —  RA 1869027", COR_WILLIAN, "1")

txt(s, "Config Bits + init_hardware()", 0.5, 0.7, 12, 0.7,
    size=28, bold=True, color=COR_WILLIAN)
divider(s, 1.42, color=COR_WILLIAN)

txt(s, "Configuração do oscilador e dos pinos I/O", 0.5, 1.5, 6.5, 0.45,
    size=15, color=GRAY)

code_cfg = [
    "#pragma config FOSC   = HS         // cristal externo 20 MHz",
    "#pragma config CPUDIV = OSC1_PLL2  // CPU = Fosc (sem divisão PLL)",
    "#pragma config WDT    = OFF        // watchdog desabilitado",
    "#pragma config PBADEN = OFF        // PORTB digital no reset",
    "#pragma config LVP    = OFF        // sem programação LVP",
]
txt_code(s, code_cfg, 0.3, 1.95, 6.4, 1.65, size=12)

code_hw = [
    "void init_hardware(void) {",
    "    ADCON1 = 0x0F;           // todos pinos digitais (AN0 config. no init_adc)",
    "",
    "    TRISB = 0xFF;            // RB0-RB5 = entradas (botões)",
    "    INTCON2bits.RBPU = 0;    // 0 = habilita pull-ups internos PORTB",
    "",
    "    TRISDbits.TRISD0 = 0;   // motor IN1 — saída",
    "    TRISDbits.TRISD1 = 0;   // motor IN2 — saída",
    "    TRISDbits.TRISD2 = 0;   // LED porta — saída",
    "    LATD = 0x00;             // tudo em zero",
    "",
    "    TRISCbits.TRISC2 = 0;   // RC2 = saída PWM (CCP1)",
    "    TRISAbits.TRISA0 = 1;   // RA0 = entrada analógica (AN0)",
    "}",
]
txt_code(s, code_hw, 0.3, 3.65, 6.4, 3.5, size=11)

# bloco de notas
box(s, 7.0, 1.95, 5.8, 5.2, color=BG_CARD)
box(s, 7.0, 1.95, 5.8, 0.07, color=COR_WILLIAN)
txt(s, "Por que FOSC = HS?", 7.15, 2.05, 5.5, 0.4,
    size=14, bold=True, color=COR_WILLIAN)
notas = [
    "HS (High Speed) = cristal 20 MHz",
    "CPUDIV = OSC1_PLL2 mantém CPU em Fosc",
    "(PLL não é ativado — não precisamos de USB)",
    "",
    "Por que ADCON1 = 0x0F?",
    "PCFG=1111 desabilita TODOS os canais",
    "analógicos. Assim PORTB e PORTD ficam",
    "100% digitais no boot.",
    "AN0 é reabilitado depois em init_adc().",
    "",
    "RBPU = 0 → habilita resistores internos",
    "de pull-up em PORTB, garantindo nível",
    "HIGH quando o botão não está pressionado.",
]
bullet_block(s, notas, 7.15, 2.5, 5.5, 4.5, size=13, icon="→", color=WHITE)

# ═══════════════════════════════════════════════════════════════
#  SLIDE 6 — WILLIAN: PWM
# ═══════════════════════════════════════════════════════════════
s = add_slide()
header_bar(s, "WILLIAN DOUGLAS  —  RA 1869027", COR_WILLIAN, "1")

txt(s, "PWM — 7 kHz e 38% Duty Cycle", 0.5, 0.7, 12, 0.7,
    size=28, bold=True, color=COR_WILLIAN)
divider(s, 1.42, color=COR_WILLIAN)

# formulas
formulas = [
    ("Fpwm  =", "Fosc ÷ (4 × Prescaler × (PR2 + 1))",
     "= 20.000.000 ÷ (4 × 4 × 179)  =  6.983 Hz  ≈  7 kHz  ✓"),
    ("DC10bit =","int( 0,38 × 4 × (PR2 + 1) )",
     "= int( 0,38 × 716 )  =  272"),
    ("CCPR1L  =","DC10bit >> 2",
     "= 272 >> 2  =  68"),
    ("DC1B1:0 =","DC10bit & 0x03",
     "= 272 & 3  =  0  → CCP1CON = 0x0C"),
    ("Duty real=","272 ÷ 716",
     "= 37,99%  ≈  38%  ✓"),
]

for i, (label, formula, result) in enumerate(formulas):
    y = 1.55 + i * 0.88
    box(s, 0.3, y, 2.0, 0.72, color=BG_CARD)
    box(s, 2.3, y, 4.5, 0.72, color=BG_CARD)
    box(s, 6.8, y, 6.1, 0.72, color=RGBColor(0x1E,0x25,0x33))
    txt(s, label,   0.4,  y+0.15, 2.0, 0.5, size=13, bold=True, color=ACCENT)
    txt(s, formula, 2.4,  y+0.15, 4.4, 0.5, size=13, color=WHITE)
    txt(s, result,  6.9,  y+0.15, 5.9, 0.5, size=13, color=GREEN, bold=True)

code_pwm = [
    "void init_pwm(void) {",
    "    OpenTimer2(TIMER_INT_OFF & T2_PS_1_4 & T2_POST_1_1); // psc 1:4",
    "    OpenPWM1(178);    // PR2 = 178 → Fpwm ≈ 7 kHz",
    "    SetDCPWM1(0);     // inicia desligado (duty = 0%)",
    "}",
    "// Ligar motor (38% duty):",
    "SetDCPWM1(272);   // duty 10-bit = 272 → 37,99% ≈ 38%",
    "// Desligar motor:",
    "SetDCPWM1(0);     // duty = 0%",
]
txt_code(s, code_pwm, 0.3, 5.9, 12.7, 1.6, size=11)

# ═══════════════════════════════════════════════════════════════
#  SLIDE 7 — KAROLINE: Timer0 ISR
# ═══════════════════════════════════════════════════════════════
s = add_slide()
header_bar(s, "KAROLINE VITÓRIA  —  RA 2578000", COR_KAROLINE, "2")

txt(s, "Timer0 — Base de Tempo de 1 ms (ISR)", 0.5, 0.7, 12, 0.7,
    size=28, bold=True, color=COR_KAROLINE)
divider(s, 1.42, color=COR_KAROLINE)

calcs = [
    ("Fosc = 20 MHz", "T_maq = 4 / 20 MHz = 200 ns"),
    ("Prescaler 1:8",  "T_tick = 8 × 200 ns = 1,6 µs por tick"),
    ("Para 1 ms:",     "1.000 µs ÷ 1,6 µs = 625 ticks"),
    ("Recarga Timer0:","65.536 − 625  =  64.911  =  0xFD8F"),
]
for i, (label, value) in enumerate(calcs):
    x = 0.3 if i < 2 else 6.8
    y = 1.55 + (i % 2) * 0.85
    box(s, x, y, 2.2, 0.72, color=COR_KAROLINE)
    box(s, x+2.2, y, 4.0, 0.72, color=BG_CARD)
    txt(s, label, x+0.1, y+0.15, 2.1, 0.5, size=13, bold=True, color=BG_DARK)
    txt(s, value, x+2.3, y+0.15, 3.9, 0.5, size=13, color=WHITE)

code_isr = [
    "// Vetor de interrupção alta prioridade (endereço 0x0008)",
    "#pragma code int_alta = 0x0008",
    "void int_alta(void) { _asm GOTO ISR_TIMER0 _endasm }",
    "#pragma code",
    "",
    "#pragma interrupt ISR_TIMER0",
    "void ISR_TIMER0(void) {",
    "    ms_tick++;          // incrementa o contador global de ms",
    "    TMR0H = 0xFD;       // recarga para próxima interrupção em 1 ms",
    "    TMR0L = 0x8F;",
    "    INTCONbits.TMR0IF = 0;  // limpa flag de overflow",
    "}",
]
txt_code(s, code_isr, 0.3, 3.3, 7.8, 3.1, size=11)

box(s, 8.3, 3.3, 4.7, 3.1, color=BG_CARD)
box(s, 8.3, 3.3, 4.7, 0.07, color=COR_KAROLINE)
txt(s, "Como funciona na prática?", 8.45, 3.4, 4.5, 0.4,
    size=14, bold=True, color=COR_KAROLINE)
notas = [
    "Timer0 gera overflow a cada 625 ticks",
    "= exatamente 1 ms",
    "",
    "ms_tick é global e volátil (volatile)",
    "→ garante leitura atualizada pelo main()",
    "",
    "Medir 2300 ms de inatividade:",
    "(ms_tick − t_inativo) >= 2300",
    "",
    "Não bloqueia o loop principal — o",
    "main() continua executando e a ISR",
    "interrompe brevemente a cada 1 ms.",
]
bullet_block(s, notas, 8.45, 3.8, 4.5, 2.5, size=12, icon="→", color=WHITE)

code_init_t = [
    "// init_timer0():  16-bit | Fosc/4 | prescaler 1:8",
    "T0CON = 0b10000010;  // TMR0ON | 16-bit | int clk | psc 1:8",
    "TMR0H = 0xFD;  TMR0L = 0x8F;  // 65536 - 625 = 64.911",
    "T0CONbits.TMR0ON = 1;  // habilita Timer0",
]
txt_code(s, code_init_t, 0.3, 6.45, 7.8, 0.95, size=11)

# ═══════════════════════════════════════════════════════════════
#  SLIDE 8 — KAROLINE: ADC + Temperatura
# ═══════════════════════════════════════════════════════════════
s = add_slide()
header_bar(s, "KAROLINE VITÓRIA  —  RA 2578000", COR_KAROLINE, "2")

txt(s, "ADC — Potenciômetro Simulando Temperatura", 0.5, 0.7, 12, 0.7,
    size=26, bold=True, color=COR_KAROLINE)
divider(s, 1.42, color=COR_KAROLINE)

code_adc = [
    "void init_adc(void) {",
    "    ADCON1bits.VCFG1 = 0; ADCON1bits.VCFG0 = 0; // Vref = VDD/VSS",
    "    // PCFG = 1110: somente AN0 analógico",
    "    ADCON1bits.PCFG3=1; ADCON1bits.PCFG2=1;",
    "    ADCON1bits.PCFG1=1; ADCON1bits.PCFG0=0;",
    "    // Clock Fosc/16 → TAD = 0,8 µs (mínimo para 20 MHz)",
    "    ADCON2bits.ADCS2=1; ADCON2bits.ADCS1=0; ADCON2bits.ADCS0=1;",
    "    // Aquisição automática: 4 TAD = 3,2 µs",
    "    ADCON2bits.ACQT2=0; ADCON2bits.ACQT1=1; ADCON2bits.ACQT0=0;",
    "    ADCON2bits.ADFM = 1;   // resultado à direita (10 bits em ADRESH:ADRESL)",
    "    ADCON0 = 0x01;         // AN0 selecionado + ADON = 1",
    "}",
    "",
    "unsigned int ler_adc(void) {",
    "    ADCON0bits.GO = 1;         // inicia conversão",
    "    while (!PIR1bits.ADIF);    // aguarda fim (~12 µs)",
    "    PIR1bits.ADIF = 0;",
    "    return ((unsigned int)ADRESH << 8) | ADRESL;  // 0–1023",
    "}",
]
txt_code(s, code_adc, 0.3, 1.55, 7.7, 4.95, size=11)

box(s, 8.2, 1.55, 4.8, 4.95, color=BG_CARD)
box(s, 8.2, 1.55, 4.8, 0.07, color=COR_KAROLINE)
txt(s, "Mapeamento de temperatura", 8.35, 1.65, 4.6, 0.45,
    size=14, bold=True, color=COR_KAROLINE)
notas = [
    "Potenciômetro em RA0 (AN0) gera 0–5V",
    "ADC 10 bits: 0 V = 0 | 5 V = 1023",
    "",
    "Mapeamento → 0 a 50 °C:",
    "temp = adc × 50 / 1023",
    "",
    "Com 1 decimal (sem float):",
    "temp_10 = adc × 500 / 1023",
    "→ ex: adc=640 → 312 → 31.2 °C",
    "",
    "Limiar de parada: 31 °C",
    "→ adc ≈ 634 (62% do potenc.)",
    "",
    "Não usamos float — operações",
    "com unsigned long int evitam",
    "lentidão e erros no C18 LITE.",
]
bullet_block(s, notas, 8.35, 2.1, 4.6, 4.3, size=12, icon="→", color=WHITE)

# Fórmula visual
box(s, 0.3, 6.6, 12.7, 0.75, color=RGBColor(0x1E,0x25,0x33))
txt(s, "temp = (unsigned long)adc_raw × 50UL / 1023UL     |     "
       "temp_10 = (unsigned long)adc_raw × 500UL / 1023UL",
    0.5, 6.68, 12.4, 0.55, size=15, color=GREEN, align=PP_ALIGN.CENTER, bold=True)

# ═══════════════════════════════════════════════════════════════
#  SLIDE 9 — ANA CAROLINE: LCD + Miyadaira
# ═══════════════════════════════════════════════════════════════
s = add_slide()
header_bar(s, "ANA CAROLINE  —  RA 2486024", COR_ANA, "3")

txt(s, "LCD 2×16 — Biblioteca Miyadaira", 0.5, 0.7, 12, 0.7,
    size=28, bold=True, color=COR_ANA)
divider(s, 1.42, color=COR_ANA)

# tabela de funções
funcoes = [
    ("lcd_inicia(c1, c2, c3)", "Inicializa o display (modo, cursor, shift)"),
    ("lcd_posicao(linha, col)", "Posiciona cursor — 1-based (linha 1 ou 2)"),
    ("lcd_escreve_dado(byte)",  "Escreve 1 caractere/símbolo no cursor atual"),
    ("imprime_string_lcd(str)", "Envia string da memória de programa para o LCD"),
    ("imprime_buffer_lcd(buf,n)","Envia n bytes de um buffer RAM para o LCD"),
    ("lcd_envia_controle(RS,RW,d,t)","Envia comando bruto (RS=0) ou dado (RS=1)"),
]

box(s, 0.3, 1.55, 5.5, 0.42, color=COR_ANA)
box(s, 5.8, 1.55, 7.1, 0.42, color=COR_ANA)
txt(s, "Função",    0.4, 1.58, 5.4, 0.38, size=13, bold=True, color=BG_DARK)
txt(s, "Descrição", 5.9, 1.58, 7.0, 0.38, size=13, bold=True, color=BG_DARK)

for i, (fn, desc) in enumerate(funcoes):
    y = 1.97 + i * 0.52
    bg = BG_CARD if i % 2 == 0 else RGBColor(0x1E,0x25,0x33)
    box(s, 0.3, y, 5.5, 0.5, color=bg)
    box(s, 5.8, y, 7.1, 0.5, color=bg)
    txt(s, fn,   0.4, y+0.1, 5.4, 0.4, size=12, color=GREEN, bold=True)
    txt(s, desc, 5.9, y+0.1, 7.0, 0.4, size=12, color=WHITE)

code_init = [
    "// Inicialização usada no projeto:",
    "lcd_inicia(0x28, 0x0F, 0x06);",
    "//         0x28 = 4 bits, 2 linhas, matriz 8x5",
    "//         0x0F = liga display + cursor piscante",
    "//         0x06 = cursor desloca para direita",
    "",
    "// Escrever na linha 1 coluna 1:",
    "lcd_posicao(1, 1);",
    "imprime_string_lcd(\"Andar: \");",
    "lcd_escreve_dado('3');",
]
txt_code(s, code_init, 0.3, 5.1, 12.7, 2.3, size=12)

# ═══════════════════════════════════════════════════════════════
#  SLIDE 10 — ANA CAROLINE: CGRAM (chars customizados)
# ═══════════════════════════════════════════════════════════════
s = add_slide()
header_bar(s, "ANA CAROLINE  —  RA 2486024", COR_ANA, "3")

txt(s, "CGRAM — Criando os Símbolos ↑ e ↓", 0.5, 0.7, 12, 0.7,
    size=28, bold=True, color=COR_ANA)
divider(s, 1.42, color=COR_ANA)

# bitmaps visuais
for ci, (nome, bm, sx) in enumerate([
    ("↑  char 0", [0x04,0x0E,0x1F,0x04,0x04,0x04,0x04,0x00], 0.3),
    ("↓  char 1", [0x04,0x04,0x04,0x04,0x1F,0x0E,0x04,0x00], 3.5),
]):
    txt(s, nome, sx, 1.5, 2.8, 0.45, size=16, bold=True, color=COR_ANA)
    for row, byte in enumerate(bm):
        for col in range(5):
            bit = (byte >> (4 - col)) & 1
            color = WHITE if bit else RGBColor(0x30,0x37,0x42)
            cell = box(s, sx + col * 0.33, 1.98 + row * 0.33, 0.3, 0.3, color=color)
        # valor hex
        txt(s, f"  0x{byte:02X}", sx + 5 * 0.33 + 0.05, 1.98 + row * 0.33,
            0.9, 0.3, size=11, color=GRAY)

txt(s, "O HD44780 suporta até 8 chars customizados na CGRAM (endereços 0x40–0x7F, 8 bytes cada).",
    0.3, 4.65, 12.7, 0.38, size=13, color=GRAY)

code_cgram = [
    "const unsigned char BITMAP_CIMA[8]  = {0x04,0x0E,0x1F,0x04,0x04,0x04,0x04,0x00};",
    "const unsigned char BITMAP_BAIXO[8] = {0x04,0x04,0x04,0x04,0x1F,0x0E,0x04,0x00};",
    "",
    "void carregar_chars(void) {",
    "    unsigned char i;",
    "    lcd_envia_controle(0, 0, 0x40, 45); // aponta CGRAM para char 0",
    "    for (i = 0; i < 8; i++) lcd_escreve_dado(BITMAP_CIMA[i]);",
    "    lcd_envia_controle(0, 0, 0x48, 45); // aponta CGRAM para char 1",
    "    for (i = 0; i < 8; i++) lcd_escreve_dado(BITMAP_BAIXO[i]);",
    "    lcd_posicao(1, 1);  // volta cursor para DDRAM",
    "}",
    "// Exibir: lcd_escreve_dado(0x00) = seta cima | lcd_escreve_dado(0x01) = seta baixo",
]
txt_code(s, code_cgram, 0.3, 5.08, 12.7, 2.35, size=10)

# ═══════════════════════════════════════════════════════════════
#  SLIDE 11 — ANA CAROLINE: Botões + Debounce
# ═══════════════════════════════════════════════════════════════
s = add_slide()
header_bar(s, "ANA CAROLINE  —  RA 2486024", COR_ANA, "3")

txt(s, "Botões — ler_botoes() com Debounce", 0.5, 0.7, 12, 0.7,
    size=28, bold=True, color=COR_ANA)
divider(s, 1.42, color=COR_ANA)

code_btn = [
    "static unsigned long t_deb[7] = {0};  // timestamp do último debounce por andar",
    "",
    "void ler_botoes(void) {",
    "    unsigned char i, portb;",
    "    portb = PORTB;   // lê o porto inteiro de uma vez (snapshot)",
    "",
    "    for (i = 1; i <= TOTAL_ANDARES; i++) {",
    "        // bit (i-1) do PORTB → botão do andar i",
    "        // pull-up ativo: botão pressionado = LOW (0)",
    "        if (!(portb & (1 << (i - 1)))) {",
    "            // debounce: só aceita se passaram ≥ 20 ms desde o último",
    "            if ((ms_tick - t_deb[i]) >= 20) {",
    "                t_deb[i] = ms_tick;",
    "                if (i != (unsigned char)andar_atual) {",
    "                    fila[i]   = 1;       // adiciona à fila",
    "                    t_inativo = ms_tick; // reinicia contador de inatividade",
    "                }",
    "            }",
    "        }",
    "    }",
    "}",
]
txt_code(s, code_btn, 0.3, 1.55, 8.1, 5.35, size=11)

box(s, 8.6, 1.55, 4.4, 5.35, color=BG_CARD)
box(s, 8.6, 1.55, 4.4, 0.07, color=COR_ANA)
txt(s, "Por que debounce?", 8.75, 1.65, 4.2, 0.45,
    size=14, bold=True, color=COR_ANA)
notas_btn = [
    "Botões mecânicos 'ricocheteiam'",
    "por 5–20 ms ao serem pressionados",
    "→ geram múltiplos pulsos rápidos.",
    "",
    "Sem debounce: um toque conta",
    "como vários pressionamentos.",
    "",
    "Solução: registra o tempo do",
    "último toque (t_deb[i]).",
    "Só aceita novo toque se",
    "(ms_tick - t_deb[i]) ≥ 20 ms.",
    "",
    "t_deb[] é static → persiste",
    "entre chamadas da função.",
    "",
    "Pull-up interno (RBPU=0):",
    "lógica invertida — LOW = on.",
    "!(portb & (1<<(i-1))) detecta",
    "o bit em zero (botão ativo).",
]
bullet_block(s, notas_btn, 8.75, 2.1, 4.2, 4.7, size=11, icon="→", color=WHITE)

# ═══════════════════════════════════════════════════════════════
#  SLIDE 12 — ERICK: Máquina de Estados — Visão Geral
# ═══════════════════════════════════════════════════════════════
s = add_slide()
header_bar(s, "ERICK RODRIGUES  —  RA 2301105", COR_ERICK, "4")

txt(s, "Máquina de Estados — Visão Geral", 0.5, 0.7, 12, 0.7,
    size=28, bold=True, color=COR_ERICK)
divider(s, 1.42, color=COR_ERICK)

estados_info = [
    ("IDLE",          COR_ERICK,    0.4,  2.2, "Motor parado\nAguarda botão\nContagem 2,3 s"),
    ("MOVENDO",       ACCENT2,      4.0,  2.2, "Motor ligado\nAvança 1 andar\na cada 3 s"),
    ("RETORNANDO",    COR_KAROLINE, 7.6,  2.2, "Igual MOVENDO\npara o andar 3\n(automático)"),
    ("PORTA_ABERTA",  COR_ANA,      4.0,  4.8, "Motor parado\nPorta aberta\nAguarda 1 s"),
    ("SUPERAQUECIDO", RED,          0.4,  4.8, "Motor parado\nTemp > 31 °C\nAguarda esfriar"),
]
for nome, cor, x, y, desc in estados_info:
    box(s, x, y, 2.8, 1.8, color=BG_CARD)
    box(s, x, y, 2.8, 0.07, color=cor)
    txt(s, nome, x+0.1, y+0.15, 2.6, 0.45, size=13, bold=True, color=cor)
    txt(s, desc, x+0.1, y+0.65, 2.6, 1.1, size=12, color=WHITE)

# setas simplificadas (texto)
setas = [
    (3.2, 3.0, "→ botão\n   pressionado"),
    (6.8, 3.0, "→ 2,3 s\n   inativo"),
    (5.4, 4.5, "↓ chegou\n   no andar"),
    (2.8, 5.8, "← temp\n   > 31°C"),
    (5.4, 6.0, "↑ 1 s\n   passado"),
    (3.2, 5.8, "↑ temp ≤ 31°C\n   retoma"),
]
for x, y, label in setas:
    txt(s, label, x, y, 1.5, 0.6, size=10, color=ACCENT, bold=True)

txt(s,
    "A cada iteração do loop principal: ler_adc() → ler_botoes() → maquina_estados() → atualizar_lcd()",
    0.3, 6.9, 12.7, 0.45, size=13, color=GRAY, align=PP_ALIGN.CENTER)

# ═══════════════════════════════════════════════════════════════
#  SLIDE 13 — ERICK: Estados IDLE + MOVENDO
# ═══════════════════════════════════════════════════════════════
s = add_slide()
header_bar(s, "ERICK RODRIGUES  —  RA 2301105", COR_ERICK, "4")

txt(s, "Estados IDLE e MOVENDO/RETORNANDO", 0.5, 0.7, 12, 0.7,
    size=26, bold=True, color=COR_ERICK)
divider(s, 1.42, color=COR_ERICK)

code_idle = [
    "case IDLE:",
    "    motor_desliga();",
    "    if (temp > TEMP_MAX) { estado = SUPERAQUECIDO; break; }",
    "",
    "    andar_destino = proximo_fila();  // algoritmo SCAN",
    "    if (andar_destino != 0) {",
    "        definir_direcao();",
    "        motor_liga();",
    "        t_viagem = t_inativo = ms_tick;",
    "        estado = MOVENDO;",
    "        break;",
    "    }",
    "    // 2,3 s sem chamada e não está no 3º → retorna",
    "    if ((ms_tick - t_inativo) >= 2300UL",
    "            && andar_atual != ANDAR_HOME) {",
    "        andar_destino = ANDAR_HOME;",
    "        definir_direcao();",
    "        motor_liga();",
    "        t_viagem = t_inativo = ms_tick;",
    "        estado = RETORNANDO;",
    "    }",
    "    break;",
]
txt_code(s, code_idle, 0.3, 1.55, 6.3, 5.35, size=10)

code_mov = [
    "case MOVENDO:",
    "case RETORNANDO:",
    "    if (temp > TEMP_MAX) {",
    "        motor_desliga();",
    "        estado = SUPERAQUECIDO;",
    "        break;",
    "    }",
    "    // avança 1 andar a cada TEMPO_POR_ANDAR ms",
    "    if ((ms_tick - t_viagem) >= TEMPO_POR_ANDAR) {",
    "        t_viagem = ms_tick;",
    "        if (direcao == SUBINDO)  andar_atual++;",
    "        if (direcao == DESCENDO) andar_atual--;",
    "",
    "        if (andar_atual == andar_destino) {",
    "            fila[andar_destino] = 0; // remove da fila",
    "            motor_desliga();",
    "            porta_abre();",
    "            t_porta = ms_tick;",
    "            estado = PORTA_ABERTA;",
    "        }",
    "    }",
    "    break;",
]
txt_code(s, code_mov, 6.8, 1.55, 6.2, 5.35, size=10)

txt(s, "proximo_fila(): algoritmo SCAN — prioriza andares na direção atual; só inverte quando não há mais chamadas nessa direção.",
    0.3, 7.0, 12.7, 0.38, size=12, color=GRAY)

# ═══════════════════════════════════════════════════════════════
#  SLIDE 14 — ERICK: PORTA_ABERTA + SUPERAQUECIDO + motor/porta
# ═══════════════════════════════════════════════════════════════
s = add_slide()
header_bar(s, "ERICK RODRIGUES  —  RA 2301105", COR_ERICK, "4")

txt(s, "Estados PORTA_ABERTA e SUPERAQUECIDO", 0.5, 0.7, 12, 0.7,
    size=26, bold=True, color=COR_ERICK)
divider(s, 1.42, color=COR_ERICK)

code_porta = [
    "case PORTA_ABERTA:",
    "    // porta permanece aberta por 1 s (1000 ms)",
    "    if ((ms_tick - t_porta) >= TEMPO_PORTA) {",
    "        porta_fecha();",
    "        t_inativo = ms_tick; // reinicia inatividade",
    "        estado = IDLE;",
    "    }",
    "    break;",
    "",
    "case SUPERAQUECIDO:",
    "    motor_desliga();",
    "    direcao = PARADO;",
    "    if (temp <= TEMP_MAX) {  // temperatura normalizou",
    "        t_inativo = ms_tick;",
    "        if (andar_destino != 0",
    "                && andar_destino != andar_atual) {",
    "            definir_direcao();",
    "            motor_liga();  // retoma viagem pendente",
    "            t_viagem = ms_tick;",
    "            estado = MOVENDO;",
    "        } else {",
    "            estado = IDLE;",
    "        }",
    "    }",
    "    break;",
]
txt_code(s, code_porta, 0.3, 1.55, 6.6, 5.7, size=10)

code_motor = [
    "void motor_liga(void) {",
    "    if (direcao == SUBINDO)",
    "        { MOTOR_IN1=1; MOTOR_IN2=0; }",
    "    else",
    "        { MOTOR_IN1=0; MOTOR_IN2=1; }",
    "    SetDCPWM1(PWM_DUTY); // 38% — ≈7 kHz",
    "}",
    "void motor_desliga(void) {",
    "    SetDCPWM1(0);",
    "    MOTOR_IN1=0; MOTOR_IN2=0;",
    "}",
    "void porta_abre(void)  { PORTA_LED=1; }",
    "void porta_fecha(void) { PORTA_LED=0; }",
]
txt_code(s, code_motor, 7.1, 1.55, 5.9, 3.0, size=11)

box(s, 7.1, 4.7, 5.9, 2.5, color=BG_CARD)
box(s, 7.1, 4.7, 5.9, 0.07, color=RED)
txt(s, "Comportamento no SUPERAQUECIDO", 7.25, 4.8, 5.7, 0.45,
    size=13, bold=True, color=RED)
notas_super = [
    "Motor para imediatamente",
    "Direção é resetada para PARADO",
    "LCD mostra 'STOP' em vez de ↑/↓",
    "Quando temp cai ≤ 31 °C:",
    "  → se tinha destino: retoma viagem",
    "  → senão: volta ao IDLE normal",
]
bullet_block(s, notas_super, 7.25, 5.3, 5.6, 1.8, size=13, icon="!", color=WHITE)

# ═══════════════════════════════════════════════════════════════
#  SLIDE 15 — ERICK: atualizar_lcd()
# ═══════════════════════════════════════════════════════════════
s = add_slide()
header_bar(s, "ERICK RODRIGUES  —  RA 2301105", COR_ERICK, "4")

txt(s, "atualizar_lcd() — Display do Elevador", 0.5, 0.7, 12, 0.7,
    size=26, bold=True, color=COR_ERICK)
divider(s, 1.42, color=COR_ERICK)

code_lcd_upd = [
    "void atualizar_lcd(void) {",
    "    unsigned char buf[18];",
    "    // temperatura em décimos de grau (sem float)",
    "    unsigned int temp_10 =",
    "        (unsigned int)((unsigned long)adc_raw * 500UL / 1023UL);",
    "",
    "    // ── Linha 1: \"Andar: X  [dir]  \" ─────────────────────",
    "    lcd_posicao(1, 1);",
    "    sprintf(buf, \"Andar: %d  \", andar_atual);",
    "    imprime_buffer_lcd(buf, 10);   // 10 chars",
    "",
    "    if (estado == SUPERAQUECIDO)       imprime_string_lcd(\"STOP  \");",
    "    else if (direcao == SUBINDO)  { lcd_escreve_dado(0x00); // ↑",
    "                                    imprime_string_lcd(\"     \"); }",
    "    else if (direcao == DESCENDO) { lcd_escreve_dado(0x01); // ↓",
    "                                    imprime_string_lcd(\"     \"); }",
    "    else                              imprime_string_lcd(\"      \");",
    "",
    "    // ── Linha 2: \"Tmp: XX.X C     \" ──────────────────────",
    "    lcd_posicao(2, 1);",
    "    sprintf(buf, \"Tmp: %u.%uC       \", temp_10/10, temp_10%10);",
    "    imprime_buffer_lcd(buf, 16);",
    "}",
]
txt_code(s, code_lcd_upd, 0.3, 1.55, 8.0, 5.7, size=10)

# display mock
box(s, 8.55, 1.7, 4.4, 0.35, color=RGBColor(0x2E,0x35,0x2E))
txt(s, "LCD — Linha 1 (normal)", 8.7, 1.72, 4.2, 0.32,
    size=11, color=GRAY)
box(s, 8.55, 2.05, 4.4, 0.65, color=RGBColor(0x1A,0x3A,0x1A))
txt(s, "Andar: 3  ↑     ", 8.65, 2.12, 4.2, 0.5,
    size=16, bold=True, color=GREEN)

box(s, 8.55, 2.85, 4.4, 0.35, color=RGBColor(0x2E,0x35,0x2E))
txt(s, "LCD — Linha 2 (temperatura)", 8.7, 2.87, 4.2, 0.32,
    size=11, color=GRAY)
box(s, 8.55, 3.2, 4.4, 0.65, color=RGBColor(0x1A,0x3A,0x1A))
txt(s, "Tmp: 28.5C      ", 8.65, 3.27, 4.2, 0.5,
    size=16, bold=True, color=GREEN)

box(s, 8.55, 4.05, 4.4, 0.35, color=RGBColor(0x2E,0x35,0x2E))
txt(s, "LCD — Linha 1 (superaquecido)", 8.7, 4.07, 4.2, 0.32,
    size=11, color=GRAY)
box(s, 8.55, 4.4, 4.4, 0.65, color=RGBColor(0x3A,0x1A,0x1A))
txt(s, "Andar: 2  STOP  ", 8.65, 4.47, 4.2, 0.5,
    size=16, bold=True, color=RED)

box(s, 8.55, 5.25, 4.4, 0.35, color=RGBColor(0x2E,0x35,0x2E))
txt(s, "LCD — Linha 2 (superaquecido)", 8.7, 5.27, 4.2, 0.32,
    size=11, color=GRAY)
box(s, 8.55, 5.6, 4.4, 0.65, color=RGBColor(0x3A,0x1A,0x1A))
txt(s, "Tmp: 34.7C      ", 8.65, 5.67, 4.2, 0.5,
    size=16, bold=True, color=RED)

txt(s, "temp_10/10 = parte inteira  |  temp_10%10 = 1 decimal  →  sem float, evitando lentidão no C18 LITE",
    0.3, 7.05, 12.7, 0.38, size=12, color=GRAY)

# ═══════════════════════════════════════════════════════════════
#  SLIDE 16 — CONCLUSÃO
# ═══════════════════════════════════════════════════════════════
s = add_slide()
box(s, 0, 0, 13.33, 0.55, color=ACCENT)
txt(s, "  CONCLUSÃO", 0.1, 0.05, 12, 0.45, size=18, bold=True, color=BG_DARK)

txt(s, "O que foi implementado", 0.5, 0.75, 12, 0.65,
    size=28, bold=True, color=WHITE)
divider(s, 1.45)

conclusoes = [
    "Timer0 (16 bits, ISR) gera base de tempo de 1 ms não-bloqueante para todas as contagens",
    "PWM no CCP1/RC2: 6.983 Hz ≈ 7 kHz com 38% duty cycle calculado analiticamente",
    "ADC em AN0 com Fosc/16 lê o potenciômetro e mapeia 0–50 °C sem usar float",
    "LCD 4 bits (Miyadaira) exibe andar, temperatura e símbolos ↑↓ criados na CGRAM",
    "Máquina de estados SCAN com 5 estados gerencia toda a lógica do elevador",
    "Debounce por tempo (20 ms) garante leitura correta dos 6 botões em PORTB",
]
bullet_block(s, conclusoes, 0.5, 1.6, 12.3, 4.0, size=17, icon="✔", color=WHITE)

# membros no final
for i, ((nome, ra, cor), xs2) in enumerate(zip(membros, [0.3, 3.6, 6.9, 10.0])):
    box(s, xs2, 5.85, 3.0, 1.3, color=BG_CARD)
    box(s, xs2, 5.85, 3.0, 0.06, color=cor)
    partes_label = ["Parte 1", "Parte 2", "Parte 3", "Parte 4"]
    txt(s, partes_label[i], xs2+0.12, 5.93, 2.8, 0.3,
        size=10, color=cor, bold=True, align=PP_ALIGN.CENTER)
    txt(s, nome, xs2+0.1, 6.25, 2.8, 0.55,
        size=11, color=WHITE, align=PP_ALIGN.CENTER)
    txt(s, ra, xs2+0.1, 6.8, 2.8, 0.3,
        size=11, color=GRAY, align=PP_ALIGN.CENTER)

txt(s, "ELE1012 — Microcontroladores  |  Prof. Miyadaira  |  UTFPR Medianeira  |  2026",
    0.5, 7.1, 12.33, 0.35, size=12, color=GRAY, align=PP_ALIGN.CENTER)

# ─── SALVAR ──────────────────────────────────────────────────
OUTPUT = r"D:\Projetos\UTFPR\microcontroladores\elevador\Central_Elevador.pptx"
prs.save(OUTPUT)
print(f"PPT gerado: {OUTPUT}")
print(f"Total de slides: {len(prs.slides)}")
